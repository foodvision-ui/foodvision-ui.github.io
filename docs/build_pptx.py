#!/usr/bin/env python3
"""
Build two PPTX files from docs/slides.html:
  1) slides_hifi.pptx   — each slide = high-res screenshot of the HTML (layout preserved 1:1)
  2) slides_native.pptx — each slide reconstructed with python-pptx text boxes/shapes

Both decks include PowerPoint-native slide transitions and entrance animations
to approximate the HTML's CSS transitions.
"""
import os
import re
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
SLIDES_HTML = HERE / "slides.html"
CAPTURE_DIR = Path("/tmp/foodvision_slides_capture")
OUT_HIFI = HERE / "slides_hifi.pptx"
OUT_NATIVE = HERE / "slides_native.pptx"

# 16:9 at 1600x900 px, output size for PowerPoint (13.333" x 7.5")
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# ───────────────────────── Capture pipeline ─────────────────────────

def make_capture_html(n_slides: int) -> Path:
    """
    Write a modified copy of slides.html that:
      - Hides the chrome (controls, counter, progress, notes)
      - Removes slide transitions (active slide is immediately visible, content is static)
      - Leaves the URL hash mechanism intact
    """
    CAPTURE_DIR.mkdir(parents=True, exist_ok=True)
    src = SLIDES_HTML.read_text(encoding="utf-8")

    override = r"""
<style id="__capture_overrides__">
  /* Hide operator chrome */
  .controls, .counter, .progress, .notes-panel { display: none !important; }
  /* Kill the sibling slide offset so nothing hangs outside the viewport */
  .slide { transition: none !important; }
  .slide:not(.active) { display: none !important; }
  .slide.active { opacity: 1 !important; transform: none !important; }
  /* Force all content animations to their "done" state */
  .slide.active .content > *,
  .slide .content > * {
    opacity: 1 !important;
    transform: none !important;
    animation: none !important;
    transition: none !important;
  }
  /* Hide scrollbars */
  html, body { overflow: hidden !important; background: #fff9f2 !important; }
</style>
"""
    # Inject just before </head>
    assert "</head>" in src
    src = src.replace("</head>", override + "</head>")

    out = CAPTURE_DIR / "slides_capture.html"
    out.write_text(src, encoding="utf-8")
    return out


def count_slides() -> int:
    return len(re.findall(r'class="slide"', SLIDES_HTML.read_text(encoding="utf-8")))


def screenshot_slide(capture_html: Path, n: int, out_png: Path) -> None:
    # Wait long enough for webfonts to arrive
    url = f"file://{capture_html}#{n}"
    cmd = [
        "google-chrome",
        "--headless=new",
        "--disable-gpu",
        "--no-sandbox",
        "--hide-scrollbars",
        "--force-device-scale-factor=1",
        "--virtual-time-budget=15000",
        "--run-all-compositor-stages-before-draw",
        f"--window-size=1600,900",
        f"--screenshot={out_png}",
        url,
    ]
    subprocess.run(cmd, check=True, capture_output=True)


def capture_all() -> list[Path]:
    n = count_slides()
    capture_html = make_capture_html(n)
    pngs: list[Path] = []
    for i in range(1, n + 1):
        png = CAPTURE_DIR / f"slide_{i:02d}.png"
        screenshot_slide(capture_html, i, png)
        pngs.append(png)
        print(f"  captured slide {i:02d} → {png.name}")
    return pngs


# ───────────────────────── Path 1: image-based PPTX ─────────────────────────

def build_hifi_pptx(pngs: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    for i, png in enumerate(pngs, start=1):
        slide = prs.slides.add_slide(blank)
        # Full-bleed image
        pic = slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
        # Add a fade transition on this slide
        _add_slide_transition(slide, kind="fade", duration_ms=600)
        # Add a fade-in entrance on the picture
        _add_fade_in_animation(slide, pic)
    prs.save(OUT_HIFI)


# ─────────────────── PPTX XML helpers for transitions/animations ───────────────────

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

def _add_slide_transition(slide, kind: str = "fade", duration_ms: int = 500) -> None:
    """Attach a <p:transition> block to the slide XML (fade by default)."""
    # Remove any pre-existing transition
    existing = slide.element.find(qn("p:transition"))
    if existing is not None:
        slide.element.remove(existing)
    xml = (
        f'<p:transition xmlns:p="{P_NS}" spd="med" dur="{duration_ms}">'
        f'  <p:{kind}/>'
        f'</p:transition>'
    )
    slide.element.append(etree.fromstring(xml))


def _add_fade_in_animation(slide, shape, delay_ms: int = 0, dur_ms: int = 600) -> None:
    """Add a simple fade-in entrance animation to `shape` when the slide starts."""
    sp_id = shape.shape_id
    # Build <p:timing> → main sequence → fade entrance
    timing_xml = f'''
<p:timing xmlns:p="{P_NS}" xmlns:a="{A_NS}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" presetID="10" presetClass="entr" presetSubtype="0"
                                    fill="hold" grpId="0" nodeType="clickEffect">
                                <p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst>
                                <p:childTnLst>
                                  <p:set>
                                    <p:cBhvr>
                                      <p:cTn id="6" dur="1" fill="hold">
                                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                      </p:cTn>
                                      <p:tgtEl><p:spTgt spid="{sp_id}"/></p:tgtEl>
                                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                                    </p:cBhvr>
                                    <p:to><p:strVal val="visible"/></p:to>
                                  </p:set>
                                  <p:anim calcmode="lin" valueType="num">
                                    <p:cBhvr additive="base">
                                      <p:cTn id="7" dur="{dur_ms}" fill="hold"/>
                                      <p:tgtEl><p:spTgt spid="{sp_id}"/></p:tgtEl>
                                      <p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst>
                                    </p:cBhvr>
                                    <p:tavLst>
                                      <p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav>
                                      <p:tav tm="100000"><p:val><p:fltVal val="1"/></p:val></p:tav>
                                    </p:tavLst>
                                  </p:anim>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
            <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
  <p:bldLst>
    <p:bldP spid="{sp_id}" grpId="0"/>
  </p:bldLst>
</p:timing>
'''
    # Replace any existing timing block on the slide
    existing = slide.element.find(qn("p:timing"))
    if existing is not None:
        slide.element.remove(existing)
    slide.element.append(etree.fromstring(timing_xml))


# ───────────────────────── Path 2: native PPTX ─────────────────────────

ORANGE = RGBColor(0xFF, 0x66, 0x33)
ORANGE_SOFT = RGBColor(0xFF, 0x9F, 0x43)
CREAM = RGBColor(0xFF, 0xF9, 0xF2)
CREAM_WARM = RGBColor(0xFF, 0xE9, 0xD6)
INK = RGBColor(0x3D, 0x29, 0x1E)
INK_SOFT = RGBColor(0x6F, 0x4A, 0x36)
MUTED = RGBColor(0x80, 0x59, 0x41)
GREEN = RGBColor(0x1F, 0x9D, 0x55)
GREEN_BG = RGBColor(0xE6, 0xF4, 0xEA)
CHIP = RGBColor(0xFF, 0xF1, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _set_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _set_no_line(shape) -> None:
    shape.line.fill.background()


def _new_slide(prs, with_bg: bool = True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if with_bg:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        _set_fill(bg, CREAM)
        _set_no_line(bg)
        # Warm gradient-ish overlay using a second, shifted rect at low-ish opacity
        warm = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(8), Inches(8))
        _set_fill(warm, CREAM_WARM)
        _set_no_line(warm)
        warm.fill.fore_color.brightness = 0.0

        orb1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.2), Inches(-0.5), Inches(3.6), Inches(3.6))
        _set_fill(orb1, ORANGE_SOFT)
        _set_no_line(orb1)

        orb2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.0), Inches(5.6), Inches(3.2), Inches(3.2))
        _set_fill(orb2, ORANGE)
        _set_no_line(orb2)
    return slide


def _txt(slide, x, y, w, h, text, size, *, bold=False, color=INK, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font="Outfit"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _runs(para, parts):
    """parts = [(text, bold, color, size)]. First run replaces any default run."""
    # Clear any existing runs
    for r in list(para.runs):
        r._r.getparent().remove(r._r)
    for text, bold, color, size in parts:
        run = para.add_run()
        run.text = text
        run.font.name = "Outfit"
        run.font.bold = bold
        run.font.size = Pt(size)
        run.font.color.rgb = color


def _chip(slide, x, y, text, *, fill=CHIP, color=ORANGE, size=12, w=None):
    # A pill-shaped chip
    width = Inches(w) if w else Inches(len(text) * 0.12 + 0.5)
    h = Inches(0.42)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), width, h)
    pill.adjustments[0] = 0.5
    _set_fill(pill, fill)
    _set_no_line(pill)
    tf = pill.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.name = "Outfit"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    return pill


def _footer(slide, label: str, right: str):
    line = slide.shapes.add_connector(1, Inches(0.6), Inches(7.0), Inches(12.733), Inches(7.0))
    line.line.color.rgb = RGBColor(0xE5, 0xD0, 0xB5)
    line.line.width = Pt(0.5)
    _txt(slide, 0.6, 7.1, 6, 0.35, "● " + label.upper(), 10, bold=True, color=MUTED)
    _txt(slide, 7.0, 7.1, 6.0, 0.35, right.upper(), 10, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)


def _card(slide, x, y, w, h, *, fill=WHITE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.adjustments[0] = 0.06
    _set_fill(card, fill)
    card.line.color.rgb = RGBColor(0xF1, 0xDE, 0xC4)
    card.line.width = Pt(0.75)
    return card


def build_native_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # ───── Slide 1 — Title
    s = _new_slide(prs)
    _chip(s, 0.83, 2.4, "EECS 449 · FINAL PROJECT", fill=CHIP, color=ORANGE, size=12, w=3.2)
    box = s.shapes.add_textbox(Inches(0.83), Inches(2.9), Inches(11), Inches(2.0))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = "FoodVision "; r1.font.name = "Outfit"; r1.font.size = Pt(96); r1.font.bold = True; r1.font.color.rgb = INK
    r2 = p.add_run(); r2.text = "AI."; r2.font.name = "Outfit"; r2.font.size = Pt(96); r2.font.bold = True; r2.font.color.rgb = ORANGE
    _txt(s, 0.83, 4.4, 11, 0.8, "Meal planning that learns from you,", 26, color=INK_SOFT)
    _txt(s, 0.83, 4.9, 11, 0.8, "not from your onboarding form.", 26, color=INK_SOFT)
    _txt(s, 0.83, 6.2, 11, 0.5, "Yizhong Zhong · April 2026", 16, color=MUTED)
    _footer(s, "FoodVision AI", "01 · Title")
    _add_slide_transition(s, "fade", 500)

    # ───── Slide 2 — The fridge problem (hook)
    s = _new_slide(prs)
    _chip(s, 0.83, 1.0, "The hook", fill=CHIP, color=ORANGE, size=12, w=1.35)
    _txt(s, 0.83, 1.8, 12, 1.2, "It's 7pm.", 84, bold=True, color=INK)
    _txt(s, 0.83, 3.2, 12, 1.2, "You open the fridge.", 84, bold=True, color=INK)
    _txt(s, 0.83, 4.6, 12, 1.2, "You don't know what to make.", 84, bold=True, color=ORANGE)
    _footer(s, "Hook", "02 · Problem")
    _add_slide_transition(s, "fade", 500)

    # ───── Slide 3 — Options today
    s = _new_slide(prs)
    _chip(s, 0.83, 0.9, "The options today", fill=CHIP, color=ORANGE, size=12, w=2.6)
    _txt(s, 0.83, 1.4, 12, 1, "Every food app, same failure modes.", 44, bold=True, color=INK)
    cards = [
        ("Google it", "Same 10 recipes everyone else sees."),
        ("Personalized apps", "40-field profile before your first meal."),
        ("The \"smart\" ones", "Never update after the first week."),
    ]
    x = 0.83
    for title, body in cards:
        _card(s, x, 3.0, 4.0, 3.0)
        _txt(s, x + 0.3, 3.25, 3.5, 0.4, "✕  " + title.upper(), 12, bold=True, color=ORANGE)
        _txt(s, x + 0.3, 3.8, 3.5, 0.6, title, 28, bold=True, color=INK)
        _txt(s, x + 0.3, 4.6, 3.5, 1.5, body, 16, color=INK_SOFT)
        x += 4.15
    _footer(s, "Failure modes", "03 · Today")
    _add_slide_transition(s, "fade", 500)

    # ───── Slide 4 — The pitch
    s = _new_slide(prs)
    _chip(s, 0.83, 1.0, "FoodVision AI", fill=CHIP, color=ORANGE, size=12, w=1.9)
    _txt(s, 0.83, 2.1, 12, 1.3, "Photo in. Recipe out.", 84, bold=True, color=INK)
    _txt(s, 0.83, 3.5, 12, 1.3, "Profile builds itself.", 84, bold=True, color=ORANGE)
    _txt(s, 0.83, 5.4, 12, 1, "Starts with zero knowledge about you.", 24, color=INK_SOFT)
    _txt(s, 0.83, 5.9, 12, 1, "Gets smarter every time you cook.", 24, color=INK_SOFT)
    _footer(s, "The pitch", "04 · Solution")
    _add_slide_transition(s, "fade", 500)

    # ───── Slide 5 — Three moments
    s = _new_slide(prs)
    _chip(s, 0.83, 0.9, "Every time you cook", fill=CHIP, color=ORANGE, size=12, w=2.55)
    _txt(s, 0.83, 1.4, 12, 1, "Three moments, every meal.", 44, bold=True, color=INK)
    steps = [
        ("01", "Photo in, recipe out", "Vision model reads your fridge and proposes a dish."),
        ("02", "Memory", "Last 14 days of meals feed the next prompt."),
        ("03", "Quiet learning", "Cooked Korean three times? Korean shows up in your profile. You never typed it."),
    ]
    x = 0.83
    for num, title, body in steps:
        _card(s, x, 3.0, 4.0, 3.3)
        _txt(s, x + 0.3, 3.25, 3.5, 0.4, f"STEP {num}", 12, bold=True, color=ORANGE)
        _txt(s, x + 0.3, 3.8, 3.5, 0.6, title, 22, bold=True, color=INK)
        _txt(s, x + 0.3, 4.7, 3.5, 2.0, body, 15, color=INK_SOFT)
        x += 4.15
    _footer(s, "How it feels", "05 · Loops")
    _add_slide_transition(s, "fade", 500)

    # ───── Slide 6 — Landing page transition
    s = _new_slide(prs)
    _chip(s, 0.83, 0.9, "Switch to the browser", fill=CHIP, color=ORANGE, size=12, w=2.9)
    _txt(s, 0.83, 2.4, 12, 1.3, "Landing page.", 84, bold=True, color=INK)
    _txt(s, 0.83, 4.0, 12, 0.8, "→  localhost:3000", 40, bold=True, color=ORANGE, font="JetBrains Mono")
    _txt(s, 0.83, 5.3, 12, 0.6, "The same \"It's 7pm\" hook the talk opened with.", 20, color=INK_SOFT)
    _txt(s, 0.83, 5.85, 12, 0.6, "One primary CTA: Start Planning →", 20, color=INK_SOFT)
    _footer(s, "Live switch", "06 · Landing")
    _add_slide_transition(s, "push", 600)

    # ───── Slide 7 — How it works (architecture)
    s = _new_slide(prs)
    _chip(s, 0.83, 0.9, "Architecture", fill=CHIP, color=ORANGE, size=12, w=1.7)
    _txt(s, 0.83, 1.4, 12, 1, "Thin browser. Jac backend. AI at the edge.", 40, bold=True, color=INK)
    # Three boxes connected by arrows
    _card(s, 0.83, 3.2, 3.6, 2.4)
    _txt(s, 0.93, 3.35, 3.4, 0.4, "BROWSER", 12, bold=True, color=ORANGE)
    _txt(s, 0.93, 3.85, 3.4, 0.6, "HTML · JS", 22, bold=True, color=INK)
    _txt(s, 0.93, 4.6, 3.4, 1.0, "Thin UI. JWT in localStorage. Sends photo + prompt.", 14, color=INK_SOFT)
    arr1 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.6), Inches(4.1), Inches(0.8), Inches(0.5))
    _set_fill(arr1, ORANGE); _set_no_line(arr1)

    _card(s, 5.5, 3.2, 3.6, 2.4)
    _txt(s, 5.6, 3.35, 3.4, 0.4, "BACKEND", 12, bold=True, color=ORANGE)
    _txt(s, 5.6, 3.85, 3.4, 0.6, "Jaseci · Jac", 22, bold=True, color=INK)
    _txt(s, 5.6, 4.6, 3.4, 1.0, "Per-user graph. Auth. Stores profile + meals.", 14, color=INK_SOFT)
    arr2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.25), Inches(4.1), Inches(0.8), Inches(0.5))
    _set_fill(arr2, ORANGE); _set_no_line(arr2)

    _card(s, 10.2, 3.2, 2.5, 2.4)
    _txt(s, 10.3, 3.35, 2.3, 0.4, "AI", 12, bold=True, color=ORANGE)
    _txt(s, 10.3, 3.85, 2.3, 0.6, "Qwen3-VL-32B", 18, bold=True, color=INK)
    _txt(s, 10.3, 4.5, 2.3, 1.0, "via OpenRouter. Vision + reasoning.", 13, color=INK_SOFT)
    _txt(s, 0.83, 6.1, 12, 0.5, "Graph-native storage means no ORM. JWT auth is built in.", 16, color=MUTED)
    _footer(s, "Stack", "07 · Architecture")
    _add_slide_transition(s, "fade", 500)

    # ───── Slide 8 — The learning loop
    s = _new_slide(prs)
    _chip(s, 0.83, 0.9, "The learning loop", fill=CHIP, color=ORANGE, size=12, w=2.4)
    _txt(s, 0.83, 1.4, 12, 1, "After each meal, call the AI a second time.", 40, bold=True, color=INK)
    # Big quote card
    qc = _card(s, 0.83, 2.9, 11.7, 1.9)
    _set_fill(qc, CREAM_WARM)
    _txt(s, 1.1, 3.1, 11.1, 0.4, "PROMPT", 11, bold=True, color=ORANGE)
    _txt(s, 1.1, 3.55, 11.1, 1.2,
         "\"Here's the profile. Here are their recent meals. Suggest additive updates based on patterns.\"",
         22, bold=True, color=INK)
    # Two rules cards
    _card(s, 0.83, 5.2, 5.7, 1.5)
    _txt(s, 1.1, 5.4, 5.4, 0.4, "RULE 01", 11, bold=True, color=ORANGE)
    _txt(s, 1.1, 5.75, 5.4, 0.5, "Additive only.", 20, bold=True, color=INK)
    _txt(s, 1.1, 6.2, 5.4, 0.5, "The AI can never delete what the user set.", 13, color=INK_SOFT)
    _card(s, 6.83, 5.2, 5.7, 1.5)
    _txt(s, 7.1, 5.4, 5.4, 0.4, "RULE 02", 11, bold=True, color=ORANGE)
    _txt(s, 7.1, 5.75, 5.4, 0.5, "Provenance.", 20, bold=True, color=INK)
    _txt(s, 7.1, 6.2, 5.4, 0.5, "Learned fields get a badge, so it's never a lie.", 13, color=INK_SOFT)
    _footer(s, "The novel part", "08 · Loop")
    _add_slide_transition(s, "fade", 500)

    # ───── Slide 9 — Demo
    s = _new_slide(prs)
    _chip(s, 0.83, 0.9, "Live Demo", fill=CHIP, color=ORANGE, size=12, w=1.6)
    _txt(s, 0.83, 1.4, 12, 1, "Watch for four things.", 44, bold=True, color=INK)
    items = [
        ("01", "Real AI call.", "No cached responses."),
        ("02", "Three Korean meals.", "No manual preferences typed."),
        ("03", "A fourth generation with no cuisine hint.", "The recommendation still leans Korean."),
        ("04", "Log out, log back in.", "Learned preferences survive."),
    ]
    y = 2.9
    for num, title, body in items:
        _txt(s, 0.83, y, 0.9, 0.7, num, 42, bold=True, color=ORANGE, font="JetBrains Mono")
        _txt(s, 1.9, y, 10.5, 0.6, title, 22, bold=True, color=INK)
        _txt(s, 1.9, y + 0.6, 10.5, 0.5, body, 16, color=INK_SOFT)
        y += 1.0
    _footer(s, "Demo in browser", "09 · Demo")
    _add_slide_transition(s, "push", 700)

    # ───── Slide 10 — What you just saw
    s = _new_slide(prs)
    _chip(s, 0.83, 0.9, "What you just saw", fill=CHIP, color=ORANGE, size=12, w=2.4)
    _txt(s, 0.83, 1.4, 12, 1, "Three numbers.", 44, bold=True, color=INK)
    stats = [
        ("4", "Real AI calls", "No cached responses."),
        ("0", "Preference fields", "Typed by hand."),
        ("1", "Learned profile", "Survived full logout."),
    ]
    x = 0.83
    for num, k, sub in stats:
        _card(s, x, 3.0, 4.0, 3.5)
        # Big number
        _txt(s, x + 0.3, 3.1, 3.5, 2.0, num, 120, bold=True, color=ORANGE, font="JetBrains Mono")
        _txt(s, x + 0.3, 5.3, 3.5, 0.5, k, 20, bold=True, color=INK)
        _txt(s, x + 0.3, 5.85, 3.5, 0.7, sub, 14, color=INK_SOFT)
        x += 4.15
    _footer(s, "Recap", "10 · Numbers")
    _add_slide_transition(s, "fade", 500)

    # ───── Slide 11 — What was hard
    s = _new_slide(prs)
    _chip(s, 0.83, 0.9, "What was hard", fill=CHIP, color=ORANGE, size=12, w=2.0)
    _txt(s, 0.83, 1.4, 12, 1, "Three real engineering stories.", 40, bold=True, color=INK)
    items = [
        ("Cross-request visibility",
         "Jac writes in one request weren't visible to queries in the next.",
         "Fix: pass the data explicitly."),
        ("AI output reliability",
         "Qwen3-VL sometimes emits broken JSON.",
         "Fix: four-layer fallback parser."),
        ("The overwrite bug",
         "First merge clobbered existing preferences.",
         "Caught in code review, not QA. Fixed with case-insensitive union-dedupe."),
    ]
    x = 0.83
    for title, body, fix in items:
        _card(s, x, 2.9, 4.0, 4.0)
        _txt(s, x + 0.3, 3.1, 3.5, 0.5, title, 18, bold=True, color=ORANGE)
        _txt(s, x + 0.3, 3.8, 3.5, 1.8, body, 14, color=INK)
        _txt(s, x + 0.3, 5.6, 3.5, 0.4, "→ FIX", 11, bold=True, color=GREEN)
        _txt(s, x + 0.3, 6.0, 3.5, 1.3, fix, 13, color=INK_SOFT)
        x += 4.15
    _footer(s, "Honest engineering", "11 · Hard parts")
    _add_slide_transition(s, "fade", 500)

    # ───── Slide 12 — What's next
    s = _new_slide(prs)
    _chip(s, 0.83, 0.9, "What's next", fill=CHIP, color=ORANGE, size=12, w=1.6)
    _txt(s, 0.83, 1.4, 12, 1, "One done, one near, one open.", 44, bold=True, color=INK)
    items = [
        ("DONE", "Reproducible demo video", "Recorded for the submission package."),
        ("FEATURE", "\"Forget\" button", "Let users correct the AI's mistakes without resetting the whole profile."),
        ("RESEARCH", "Window size", "Is 14 days right? Longer windows drown out recent taste changes."),
    ]
    x = 0.83
    for tag, title, body in items:
        _card(s, x, 2.9, 4.0, 3.5)
        _txt(s, x + 0.3, 3.1, 3.5, 0.4, tag, 11, bold=True, color=ORANGE)
        _txt(s, x + 0.3, 3.6, 3.5, 0.7, title, 22, bold=True, color=INK)
        _txt(s, x + 0.3, 4.5, 3.5, 1.9, body, 14, color=INK_SOFT)
        x += 4.15
    _footer(s, "Roadmap", "12 · What's next")
    _add_slide_transition(s, "fade", 500)

    # ───── Slide 13 — Thanks
    s = _new_slide(prs)
    _chip(s, 0.83, 0.9, "Thanks for listening", fill=CHIP, color=ORANGE, size=12, w=2.7)
    _txt(s, 0.83, 1.7, 12, 1, "Questions?", 96, bold=True, color=INK)
    # Three meta rows
    rows = [
        ("CODE", "github.com/ZhongYeah1/449proj"),
        ("STACK", "Jaseci · Jac · Qwen3-VL-32B · OpenRouter"),
        ("CONTACT", "yeszhong@umich.edu"),
    ]
    y = 4.2
    for k, v in rows:
        _txt(s, 0.83, y, 2.0, 0.5, k, 13, bold=True, color=ORANGE)
        _txt(s, 3.0, y, 10.0, 0.5, v, 22, bold=True, color=INK, font="JetBrains Mono")
        y += 0.7
    _footer(s, "FoodVision AI", "13 · Thanks")
    _add_slide_transition(s, "fade", 500)

    prs.save(OUT_NATIVE)


# ───────────────────────── main ─────────────────────────

def main() -> None:
    print("[1/3] Capturing HTML slides → PNG (headless Chrome)…")
    pngs = capture_all()
    print(f"       {len(pngs)} slides captured.")

    print("[2/3] Building image-based PPTX…")
    build_hifi_pptx(pngs)
    print(f"       wrote {OUT_HIFI}")

    print("[3/3] Building native PPTX…")
    build_native_pptx()
    print(f"       wrote {OUT_NATIVE}")


if __name__ == "__main__":
    main()
